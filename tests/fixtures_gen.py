"""测试夹具生成：用 python-pptx 造正文+备注 pptx；手工注入 SmartArt / 反转页序。"""
from __future__ import annotations

import shutil
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
DGM_NS = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
RELS_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
REL_DIAGRAM_DATA = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData"
)


def make_pptx(path: str | Path, slides: list[dict]) -> str:
    """slides: [{"body": str, "notes": str|None}]。返回路径。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for s in slides:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tb.text_frame.text = s.get("body", "")
        notes = s.get("notes")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
    prs.save(str(path))
    return str(path)


def reverse_slide_order(path: str | Path) -> str:
    """反转 presentation.xml 的 sldIdLst 顺序（slide 文件保持不变）。
    用于验证解析器按放映顺序、而非文件名顺序读页。"""
    path = str(path)
    with zipfile.ZipFile(path) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}
    root = etree.fromstring(data["ppt/presentation.xml"])
    lst = root.find(f"{{{P_NS}}}sldIdLst")
    children = list(lst)
    for c in children:
        lst.remove(c)
    for c in reversed(children):
        lst.append(c)
    data["ppt/presentation.xml"] = etree.tostring(
        root, xml_declaration=True, encoding="UTF-8", standalone=True
    )
    tmp = path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for n in names:
            zout.writestr(n, data[n])
    shutil.move(tmp, path)
    return path


def inject_smartart(path: str | Path, slide_no: int, text: str) -> str:
    """给第 slide_no 页注入一个含 text 的 SmartArt diagram data part + 关系。
    （slide_no 对应 ppt/slides/slideN.xml，未重排前 N==放映序）"""
    path = str(path)
    diagram_part = f"ppt/diagrams/data{slide_no}.xml"
    slide_part = f"ppt/slides/slide{slide_no}.xml"
    rels_part = f"ppt/slides/_rels/slide{slide_no}.xml.rels"

    diagram_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<dgm:dataModel xmlns:dgm="{DGM_NS}" xmlns:a="{A_NS}">'
        "<dgm:ptLst><dgm:pt><dgm:t><a:bodyPr/><a:p><a:r>"
        f"<a:t>{text}</a:t></a:r></a:p></dgm:t></dgm:pt></dgm:ptLst>"
        "</dgm:dataModel>"
    ).encode("utf-8")

    with zipfile.ZipFile(path) as zin:
        names = zin.namelist()
        data = {n: zin.read(n) for n in names}

    # 更新 slide rels：追加一个 diagramData 关系
    if rels_part in data:
        rroot = etree.fromstring(data[rels_part])
    else:
        rroot = etree.fromstring(
            f'<Relationships xmlns="{RELS_NS}"></Relationships>'.encode()
        )
    existing_ids = [r.get("Id") for r in rroot]
    n = 1
    while f"rId{n}" in existing_ids:
        n += 1
    new_id = f"rId{n}"
    rel = etree.SubElement(rroot, f"{{{RELS_NS}}}Relationship")
    rel.set("Id", new_id)
    rel.set("Type", REL_DIAGRAM_DATA)
    rel.set("Target", f"../diagrams/data{slide_no}.xml")
    data[rels_part] = etree.tostring(rroot, xml_declaration=True, encoding="UTF-8", standalone=True)
    data[diagram_part] = diagram_xml

    tmp = path + ".tmp"
    with zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for n2, d in data.items():
            zout.writestr(n2, d)
    shutil.move(tmp, path)
    return path


def make_docx(path: str | Path, paragraphs: list) -> str:
    """造最小 docx。paragraphs 每项：str=单 run 段落，list[str]=同段多 run。"""
    from xml.sax.saxutils import escape
    W = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    def runs(p) -> str:
        items = p if isinstance(p, list) else [p]
        return "".join(
            f'<w:r><w:t xml:space="preserve">{escape(t)}</w:t></w:r>' for t in items
        )

    body = "".join(f"<w:p>{runs(p)}</w:p>" for p in paragraphs)
    document = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<w:document xmlns:w="{W}"><w:body>{body}</w:body></w:document>'
    )
    ct = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>'
        "</Relationships>"
    )
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", ct)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document)
    return str(path)


def make_xlsx(path: str | Path, sheets: list) -> str:
    """造最小 xlsx（仅解析器所需 sharedStrings + 各 sheet）。sheets: list[list[str]]，每个内层列表=一张表的单元格文本。"""
    from xml.sax.saxutils import escape
    SS = "http://schemas.openxmlformats.org/spreadsheetml/2006/main"
    all_cells = [c for sheet in sheets for c in sheet]
    sis = "".join(f"<si><t>{escape(c)}</t></si>" for c in all_cells)
    shared = (
        f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        f'<sst xmlns="{SS}" count="{len(all_cells)}" uniqueCount="{len(all_cells)}">{sis}</sst>'
    )
    parts = {"xl/sharedStrings.xml": shared}
    base = 0
    for idx, sheet in enumerate(sheets, start=1):
        cells = "".join(f'<c r="A{j + 1}" t="s"><v>{base + j}</v></c>' for j in range(len(sheet)))
        parts[f"xl/worksheets/sheet{idx}.xml"] = (
            f'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            f'<worksheet xmlns="{SS}"><sheetData><row>{cells}</row></sheetData></worksheet>'
        )
        base += len(sheet)
    with zipfile.ZipFile(str(path), "w", zipfile.ZIP_DEFLATED) as z:
        for n, d in parts.items():
            z.writestr(n, d)
    return str(path)


def make_pdf(path: str | Path, pages_text: list) -> str:
    """造最小多页 PDF（ASCII 文本，可被 pypdf 抽取；偏移量程序计算保证 xref 正确）。

    pages_text: list[str]，每项一页文本。文本请用 ASCII（含 () \\ 会破坏 PDF 串语法）。
    """
    n_pages = len(pages_text)
    kids = " ".join(f"{3 + 2 * i} 0 R" for i in range(n_pages))
    font_obj = 3 + 2 * n_pages  # 字体对象号排在所有页/内容之后
    objs: list[bytes] = []
    objs.append(b"<</Type/Catalog/Pages 2 0 R>>")  # 1
    objs.append(f"<</Type/Pages/Kids[{kids}]/Count {n_pages}>>".encode())  # 2
    for i, text in enumerate(pages_text):  # 每页 2 个对象：page + content
        page_obj = 3 + 2 * i
        content_obj = page_obj + 1
        objs.append(
            f"<</Type/Page/Parent 2 0 R/MediaBox[0 0 612 792]/Contents {content_obj} 0 R"
            f"/Resources<</Font<</F1 {font_obj} 0 R>>>>>>".encode()
        )
        stream = b"BT /F1 24 Tf 72 700 Td (" + text.encode("latin-1") + b") Tj ET"
        objs.append(b"<</Length " + str(len(stream)).encode() + b">>\nstream\n" + stream + b"\nendstream")
    objs.append(b"<</Type/Font/Subtype/Type1/BaseFont/Helvetica>>")  # font

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for i, obj in enumerate(objs, start=1):
        offsets.append(len(out))
        out += f"{i} 0 obj\n".encode() + obj + b"\nendobj\n"
    xref_pos = len(out)
    out += f"xref\n0 {len(objs) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for off in offsets:
        out += f"{off:010d} 00000 n \n".encode()
    out += f"trailer\n<</Size {len(objs) + 1}/Root 1 0 R>>\nstartxref\n{xref_pos}\n%%EOF".encode()
    Path(path).write_bytes(bytes(out))
    return str(path)


def write_encrypted_stub(path: str | Path) -> str:
    """写一个 OLE 复合文档头的文件，模拟加密 pptx。"""
    Path(path).write_bytes(b"\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1" + b"\x00" * 512)
    return str(path)


def write_corrupt_stub(path: str | Path) -> str:
    Path(path).write_bytes(b"this is not a valid zip or pptx file at all")
    return str(path)
