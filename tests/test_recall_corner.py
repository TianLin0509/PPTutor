"""基础召回 corner case 验收（TDD）：PPT 原文里任意连续片段、多词组合,搜它必命中。

铁律:原文有的必召回(召回) + 不相邻的不误中(精度) + 字形归一 + 多词 + 性能<3s。
"""
from __future__ import annotations

import time

import fixtures_gen as fx
from pptx import Presentation
from pptx.util import Inches

from pptx_finder import db, indexer, search


def _build(tmp_path, files):
    """files: {filename: [page_body, ...]}（单 run）。建索引返回 conn。"""
    docs = tmp_path / "d"
    docs.mkdir(exist_ok=True)
    for fn, bodies in files.items():
        fx.make_pptx(docs / fn, [{"body": b} for b in bodies])
    conn = db.connect(tmp_path / "i.db")
    db.init_db(conn)
    indexer.update_index(conn, [str(docs)], workers=1)
    return conn


def _make_multirun(path, page_runs):
    """page_runs: [[run文本,...], ...] 每页一段落多 run（模拟 PowerPoint 拆 run）。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for runs in page_runs:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        para = tb.text_frame.paragraphs[0]
        for t in runs:
            r = para.add_run()
            r.text = t
    prs.save(str(path))
    return str(path)


def _names(res):
    return [r.name for r in res]


# ① 跨词子串（核心）
def test_01_cross_word_substring(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["小明硕士毕业典礼"]})
    assert "a.pptx" in _names(search.search(conn, "明硕"))


# ② run 截断修复：同段落被拆成两 run,仍能搜到跨 run 的「明硕」
def test_02_run_split(tmp_path):
    docs = tmp_path / "d"
    docs.mkdir()
    _make_multirun(docs / "a.pptx", [["小明", "硕士毕业"]])
    conn = db.connect(tmp_path / "i.db")
    db.init_db(conn)
    indexer.update_index(conn, [str(docs)], workers=1)
    assert "a.pptx" in _names(search.search(conn, "明硕"))


# ③ 长词中间片段
def test_03_long_word_inner(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["中华人民共和国成立"]})
    assert "a.pptx" in _names(search.search(conn, "人民"))


# ④ 多词·同页 AND
def test_04_multiword_same_page(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["小明硕士的AI研究报告"]})
    assert "a.pptx" in _names(search.search(conn, "明硕 AI"))


# ⑤ 多词·只认同一页（1A 收紧）：明硕/AI 分散在不同页 → 不再算命中
def test_05_multiword_cross_page_no_match(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["明硕方案介绍", "中间无关页", "AI落地总结"]})
    assert "a.pptx" not in _names(search.search(conn, "明硕 AI"))


# ⑥ 精度：明、硕被标点隔开,搜「明硕」不该误中
def test_06_precision_not_adjacent(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["他很聪明，硕果累累"]})
    assert "a.pptx" not in _names(search.search(conn, "明硕"))


# ⑦ 全角→半角
def test_07_fullwidth(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["采用ＡＩ技术方案"]})  # 全角 ＡＩ
    assert "a.pptx" in _names(search.search(conn, "AI"))


# ⑧ 繁体→简体
def test_08_traditional(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["軟件開發流程"]})
    assert "a.pptx" in _names(search.search(conn, "软件"))


# ⑨ 大小写
def test_09_case_insensitive(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["基于GPT的方案设计"]})
    assert "a.pptx" in _names(search.search(conn, "gpt"))


# ⑩ 数字/型号：词与子串都要中
def test_10_number_model(tmp_path):
    conn = _build(tmp_path, {"a.pptx": ["昇腾910处理器规格"]})
    assert "a.pptx" in _names(search.search(conn, "昇腾"))
    assert "a.pptx" in _names(search.search(conn, "910"))


# ⑪ 性能铁律：适中规模库,搜索 < 3 秒
def test_11_performance_under_3s(tmp_path):
    files = {
        f"f{i}.pptx": [f"第{j}页 项目{i} 通用词汇 昇腾算力集群 方案{i}_{j}" for j in range(5)]
        for i in range(40)
    }  # 40 文件 × 5 页 = 200 页
    conn = _build(tmp_path, files)
    t0 = time.time()
    search.search(conn, "昇腾 算力")
    assert time.time() - t0 < 3.0
