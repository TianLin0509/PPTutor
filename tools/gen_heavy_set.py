"""生成"重内容"PPTX 用于复现索引卡顿：多文件 × 多页 × 每页大段中文（压 tokenize/lxml CPU）。

用法：uv run python tools/gen_heavy_set.py <out_dir> [n] [pages]
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

_CJK = ("人工智能大模型算力昇腾华为研发投入战略规划市场分析技术路线产品方案客户价值"
        "数字化转型云计算边缘智能算法优化系统架构性能指标商业模式生态合作")


def gen(out_dir: str, n: int = 40, pages: int = 15) -> None:
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    body = (_CJK * 12)[:480]  # 每页约 480 字，逼真且解析有 CPU 量
    for i in range(n):
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for pg in range(pages):
            s = prs.slides.add_slide(blank)
            tb = s.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(9), Inches(6.5))
            tb.text_frame.text = f"第{pg}页 项目{i} 章节标题 " + body
        prs.save(str(out / f"heavy_{i:03d}.pptx"))
    print(f"OK: {n} files x {pages} pages -> {out}")


if __name__ == "__main__":
    gen(sys.argv[1] if len(sys.argv) > 1 else ".selftest/heavy",
        int(sys.argv[2]) if len(sys.argv) > 2 else 40,
        int(sys.argv[3]) if len(sys.argv) > 3 else 15)
