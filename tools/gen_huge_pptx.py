"""生成"巨型"PPTX：多页 + 每页大段文字（+ 内嵌图片膨胀体积），模拟真实大汇报稿。
用于复现 快照/恢复/导出 在主线程上的卡顿。

用法：uv run python tools/gen_huge_pptx.py <out_dir> [n] [slides] [with_img]
"""
from __future__ import annotations

import io
import os
import struct
import sys
import zlib
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

_CJK = ("人工智能大模型算力昇腾华为研发投入战略规划市场分析技术路线产品方案客户价值"
        "数字化转型云计算边缘智能算法优化系统架构性能指标商业模式生态合作平台能力")


def _png_blob(kb: int, seed: int = 0) -> bytes:
    """造一张可被 PowerPoint 接受的 PNG（os.urandom 随机像素压不动→膨胀体积、天然唯一不去重）。"""
    w = h = 256
    rnd = os.urandom(w * h * 3)
    raw = bytearray()
    for y in range(h):
        raw.append(0)  # filter byte
        raw += rnd[y * w * 3:(y + 1) * w * 3]
    comp = zlib.compress(bytes(raw), 1)

    def chunk(typ: bytes, data: bytes) -> bytes:
        return (struct.pack(">I", len(data)) + typ + data
                + struct.pack(">I", zlib.crc32(typ + data) & 0xFFFFFFFF))

    ihdr = struct.pack(">IIBBBBB", w, h, 8, 2, 0, 0, 0)
    png = b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr) + chunk(b"IDAT", comp) + chunk(b"IEND", b"")
    return png


def gen(out_dir: str, n: int = 3, slides: int = 300, with_img: bool = True) -> None:
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    body = (_CJK * 60)[:3000]  # 每页约 3000 字
    for i in range(n):
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for pg in range(slides):
            s = prs.slides.add_slide(blank)
            tb = s.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(9.4), Inches(6.8))
            tb.text_frame.text = f"第{pg}页 项目{i} 通用关键词 章节 " + body
            if with_img:  # 每页一张唯一图片（不去重）→ 膨胀到真实大稿体积
                s.shapes.add_picture(io.BytesIO(_png_blob(200, seed=i * 100000 + pg)),
                                     Inches(0.3), Inches(0.3), Inches(2), Inches(2))
        path = out / f"huge_{i:02d}.pptx"
        prs.save(str(path))
        print(f"  {path.name}: {slides} slides, {path.stat().st_size // 1024} KB")
    print(f"OK: {n} huge files -> {out}")


if __name__ == "__main__":
    gen(sys.argv[1] if len(sys.argv) > 1 else ".selftest/huge",
        int(sys.argv[2]) if len(sys.argv) > 2 else 3,
        int(sys.argv[3]) if len(sys.argv) > 3 else 300,
        (sys.argv[4] if len(sys.argv) > 4 else "1") != "0")
