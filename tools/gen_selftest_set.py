"""生成打包自检用的测试集（10 pptx + 4 非 pptx，名字 + 内容与 selftest.CASES 一一对应）。

用法：uv run python tools/gen_selftest_set.py <out_dir>
内容覆盖：跨词子串 / run 截断 / 长词 / 多词同页 / 多词跨页 / 精度 / 全角 / 繁简 / 大小写 / 数字
        + docx / xlsx / txt / pdf 多文档内容搜索（v1.0.0，证明 frozen 下 pypdf 等解析器带齐）。
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

# 复用测试夹具生成器造最小 docx/xlsx/pdf（pdf 的 xref 偏移由它精确计算）
sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "tests"))
import fixtures_gen as fx  # noqa: E402


def _single(path: Path, bodies: list[str]) -> None:
    """每页一个文本框、单 run。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for b in bodies:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        tb.text_frame.text = b
    prs.save(str(path))


def _multirun(path: Path, page_runs: list[list[str]]) -> None:
    """每页一段落、多 run（模拟 PowerPoint 把一句话拆成多个 <a:t>）。"""
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for runs in page_runs:
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
        para = tb.text_frame.paragraphs[0]
        for t in runs:
            para.add_run().text = t
    prs.save(str(path))


def main(out_dir: str) -> None:
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)
    _single(out / "01_cross.pptx", ["小明硕士毕业典礼"])
    _multirun(out / "02_runsplit.pptx", [["小明", "硕士毕业"]])
    _single(out / "03_longword.pptx", ["中华人民共和国成立"])
    _single(out / "04_samepage.pptx", ["小明硕士的AI研究报告"])
    _single(out / "05_crosspage.pptx", ["明硕方案介绍", "中间无关页", "AI落地总结"])
    _single(out / "06_precision.pptx", ["他很聪明，硕果累累"])
    _single(out / "07_fullwidth.pptx", ["采用ＡＩ技术方案"])  # 全角 ＡＩ
    _single(out / "08_traditional.pptx", ["軟件開發流程"])     # 繁体
    _single(out / "09_case.pptx", ["基于GPT的方案设计"])
    _single(out / "10_number.pptx", ["昇腾910处理器规格"])
    # —— v1.0.0 多文档内容搜索探针（与 selftest.CASES 11~14 对应）——
    fx.make_docx(out / "11_docx.docx", ["文档专属词 抓手落地闭环"])
    fx.make_xlsx(out / "12_xlsx.xlsx", [["表格专属词 营收同比增长"]])
    (out / "13_txt.txt").write_text("文本专属词 现金流充裕", encoding="utf-8")
    fx.make_pdf(out / "14_pdf.pdf", ["PdfSampleKeyword revenue alpha"])  # 纯 ASCII，pypdf 抽取
    print(f"OK: 14 files (10 pptx + docx/xlsx/txt/pdf) -> {out}")


if __name__ == "__main__":
    main(sys.argv[1] if len(sys.argv) > 1 else ".selftest/set")
